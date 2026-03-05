"""
Generate a styled PPTX from BugFixer2.0-Offsitev0.1-12PointFramework.html
Retains colors, fonts, styling, background colors, highlights, tables, and layout.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
from bs4 import BeautifulSoup

BASE_DIR = Path(r"c:\Users\kumar.gn\PycharmProjects\Testproject")
HTML_FILE = BASE_DIR / "BugFixer2.0-Offsitev0.1-12PointFramework.html"
OUTPUT_FILE = BASE_DIR / "BugFixer2.0-Offsitev0.1-12PointFramework.pptx"
ASSETS_DIR = BASE_DIR / "_pptx_assets"
ASSETS_DIR.mkdir(exist_ok=True)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors from HTML CSS
BLUE_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)
BLUE_DARK = RGBColor(0x1E, 0x3A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
LIGHT_BLUE_BG = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_BLUE_BG2 = RGBColor(0xBF, 0xDB, 0xFE)
YELLOW_BG = RGBColor(0xFE, 0xF0, 0x8A)
YELLOW_BORDER = RGBColor(0xCA, 0x8A, 0x04)
BROWN_TEXT = RGBColor(0x78, 0x35, 0x0F)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_BORDER = RGBColor(0x16, 0xA3, 0x4A)
GREEN_DARK = RGBColor(0x16, 0x65, 0x34)
PLACEHOLDER_BG = RGBColor(0xFE, 0xF3, 0xC7)
PLACEHOLDER_BORDER = RGBColor(0xD9, 0x77, 0x06)
PLACEHOLDER_TEXT = RGBColor(0x92, 0x40, 0x0E)
PLACEHOLDER_EM = RGBColor(0xB4, 0x53, 0x09)
GRAY_BG = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
TABLE_EVEN_ROW = RGBColor(0xF9, 0xFA, 0xFB)
SUBTITLE_COLOR = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_CELL_BG = RGBColor(0xDC, 0xFC, 0xE7)

FONT_NAME = "Segoe UI"

# ────────────────────── Image helpers ──────────────────────

def make_gradient_image(width, height, color1, color2, filename):
    """Create a linear gradient image from color1 to color2 (diagonal 135°)."""
    img = Image.new("RGB", (width, height))
    for y in range(height):
        for x in range(width):
            t = (x / width + y / height) / 2.0
            r = int(color1[0] * (1 - t) + color2[0] * t)
            g = int(color1[1] * (1 - t) + color2[1] * t)
            b = int(color1[2] * (1 - t) + color2[2] * t)
            img.putpixel((x, y), (r, g, b))
    path = ASSETS_DIR / filename
    img.save(str(path))
    return str(path)


def make_solid_image(width, height, color, filename):
    img = Image.new("RGB", (width, height), color)
    path = ASSETS_DIR / filename
    img.save(str(path))
    return str(path)


# Generate background images (small size, will be stretched)
print("Generating background images...")
title_bg_path = make_gradient_image(400, 225, (0x1E, 0x40, 0xAF), (0x1E, 0x3A, 0x8A), "title_bg.png")
highlight_bg_path = make_gradient_image(400, 60, (0xDB, 0xEA, 0xFE), (0xBF, 0xDB, 0xFE), "highlight_bg.png")
warning_bg_path = make_gradient_image(400, 60, (0xFE, 0xF0, 0x8A), (0xFD, 0xE0, 0x47), "warning_bg.png")
success_bg_path = make_gradient_image(400, 60, (0xDC, 0xFC, 0xE7), (0xBB, 0xF7, 0xD0), "success_bg.png")
white_bg_path = make_solid_image(400, 225, (0xFF, 0xFF, 0xFF), "white_bg.png")

# ────────────────────── PPTX helpers ──────────────────────

def set_slide_bg_image(slide, image_path):
    """Set a slide background to an image."""
    from pptx.oxml.ns import qn
    bg = slide.background
    fill = bg.fill
    fill.background()
    # Use picture fill
    blipFill = slide.background._element
    # Alternative: add image as full-slide picture behind everything
    slide.shapes.add_picture(image_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)


def add_bg_picture(slide, image_path):
    """Add a background image covering the full slide."""
    pic = slide.shapes.add_picture(image_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    # Move to back
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_textbox(slide, left, top, width, height, text="", font_size=14,
                font_color=DARK_TEXT, bold=False, italic=False, alignment=PP_ALIGN.LEFT,
                font_name=FONT_NAME):
    """Add a text box and return the text frame for further manipulation."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text="", font_size=14, font_color=DARK_TEXT, bold=False,
                  italic=False, alignment=PP_ALIGN.LEFT, space_before=Pt(4),
                  space_after=Pt(2), font_name=FONT_NAME, level=0):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def add_rich_paragraph(tf, runs, font_size=14, alignment=PP_ALIGN.LEFT,
                       space_before=Pt(4), space_after=Pt(2), level=0):
    """Add a paragraph with mixed bold/normal runs.
    runs is a list of (text, bold, color) tuples."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    for i, (text, bold, color) in enumerate(runs):
        if i == 0:
            run = p.runs[0] if p.runs else p.add_run()
        else:
            run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.name = FONT_NAME
    return p


def add_colored_rect(slide, left, top, width, height, fill_color, border_color=None, border_width=None):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(slide, left, top, size, fill_color, text="", font_size=24, font_color=WHITE):
    """Add a circle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    shape.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    # Vertical center
    from pptx.oxml.ns import qn
    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    bodyPr.set('anchor', 'ctr')
    return shape


# ────────────────────── Parse HTML ──────────────────────

print("Parsing HTML...")
with open(HTML_FILE, "r", encoding="utf-8") as f:
    soup = BeautifulSoup(f.read(), "html.parser")

slides_html = soup.select(".slide")


# ────────────────────── Helper to extract structured content ──────────────────────

def extract_text_from_element(el):
    """Get cleaned text from an element, preserving structure."""
    return el.get_text(strip=True) if el else ""


def parse_bullet_items(ul_el):
    """Parse bullet list items, returning list of (text, is_bold_prefix, bold_text, rest_text)."""
    items = []
    if not ul_el:
        return items
    for li in ul_el.find_all("li", recursive=False):
        strong = li.find("strong")
        if strong:
            bold_part = strong.get_text(strip=True)
            # Get rest of text after strong
            rest = li.get_text(strip=True)[len(bold_part):]
            items.append((bold_part, rest))
        else:
            # Check for nested bullets
            nested_ul = li.find("ul", class_="bullet-list")
            if nested_ul:
                text_parts = []
                for child in li.children:
                    if hasattr(child, 'name') and child.name == 'ul':
                        break
                    if hasattr(child, 'get_text'):
                        text_parts.append(child.get_text(strip=True))
                    elif isinstance(child, str) and child.strip():
                        text_parts.append(child.strip())
                main_text = " ".join(text_parts)
                items.append((None, main_text))
                # Add nested items
                for nested_li in nested_ul.find_all("li", recursive=False):
                    nested_text = nested_li.get_text(strip=True)
                    items.append(("__nested__", nested_text))
            else:
                items.append((None, li.get_text(strip=True)))
    return items


# ────────────────────── Build PPTX ──────────────────────

print("Building PPTX...")
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Use blank layout
blank_layout = prs.slide_layouts[6]

# Margins
LEFT_MARGIN = Inches(0.8)
TOP_MARGIN = Inches(0.3)
CONTENT_WIDTH = SLIDE_W - Inches(1.6)
HEADER_HEIGHT = Inches(1.1)
BODY_TOP = Inches(1.5)


def build_title_slide(slide_html, prs):
    """Build a title slide (blue gradient background, centered white text)."""
    slide = prs.slides.add_slide(blank_layout)
    add_bg_picture(slide, title_bg_path)

    # Get title and subtitles
    title_el = slide_html.select_one(".slide-title")
    paragraphs = slide_html.select("p")

    title_text = title_el.get_text(strip=True) if title_el else ""
    subtitles = [p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)]

    # Title
    tf = add_textbox(slide, LEFT_MARGIN, Inches(2.0), CONTENT_WIDTH, Inches(1.5),
                     title_text, font_size=48, font_color=WHITE, bold=True,
                     alignment=PP_ALIGN.CENTER)

    # Subtitles
    for st in subtitles:
        add_paragraph(tf, st, font_size=20, font_color=SUBTITLE_COLOR,
                      alignment=PP_ALIGN.CENTER, space_before=Pt(12), space_after=Pt(6))

    return slide


def build_content_slide(slide_html, prs):
    """Build a content slide with header bar, number circle, and content blocks."""
    slide = prs.slides.add_slide(blank_layout)

    # White background
    add_bg_picture(slide, white_bg_path)

    # Extract header info
    number_el = slide_html.select_one(".slide-number")
    title_el = slide_html.select_one(".slide-title")
    slide_num = number_el.get_text(strip=True) if number_el else ""
    slide_title = title_el.get_text(strip=True) if title_el else ""

    # Blue header line (bottom border of header)
    add_colored_rect(slide, LEFT_MARGIN, Inches(1.15), CONTENT_WIDTH, Pt(3), BLUE_PRIMARY)

    # Slide number circle
    circle_size = Inches(0.6)
    add_circle(slide, LEFT_MARGIN, Inches(0.35), circle_size, BLUE_PRIMARY,
               text=slide_num, font_size=20, font_color=WHITE)

    # Title text
    add_textbox(slide, Inches(1.6), Inches(0.3), CONTENT_WIDTH - Inches(0.8), Inches(0.8),
                slide_title, font_size=22, font_color=BLUE_PRIMARY, bold=True)

    # Now process content blocks
    content_div = slide_html.select_one(".content")
    if not content_div:
        return slide

    y_pos = Inches(1.4)
    line_spacing = Pt(2)

    # Iterate through child elements of content
    children = list(content_div.children)
    i = 0
    while i < len(children):
        el = children[i]
        if not hasattr(el, 'name') or el.name is None:
            i += 1
            continue

        classes = el.get("class", [])

        # Section Heading
        if "section-heading" in classes:
            heading_text = el.get_text(strip=True)
            tf = add_textbox(slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, Inches(0.35),
                             heading_text, font_size=15, font_color=BLUE_PRIMARY, bold=True)
            y_pos += Inches(0.38)

        # Bullet list (standalone)
        elif el.name == "ul" and "bullet-list" in classes:
            items = parse_bullet_items(el)
            tf = add_textbox(slide, Inches(1.2), y_pos, CONTENT_WIDTH - Inches(0.4), Inches(0.3),
                             "", font_size=11, font_color=DARK_TEXT)
            # Remove the empty first paragraph text
            tf.paragraphs[0].text = ""
            first = True
            for bold_part, rest_text in items:
                if bold_part == "__nested__":
                    p = add_paragraph(tf, f"    → {rest_text}", font_size=10,
                                      font_color=DARK_TEXT, space_before=Pt(1), space_after=Pt(1))
                elif bold_part:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.space_before = Pt(3)
                    p.space_after = Pt(1)
                    # Add bullet marker + bold + rest
                    run1 = p.add_run()
                    run1.text = "▸ "
                    run1.font.size = Pt(11)
                    run1.font.color.rgb = BLUE_PRIMARY
                    run1.font.bold = True
                    run1.font.name = FONT_NAME
                    run2 = p.add_run()
                    run2.text = bold_part
                    run2.font.size = Pt(11)
                    run2.font.color.rgb = DARK_TEXT
                    run2.font.bold = True
                    run2.font.name = FONT_NAME
                    if rest_text:
                        run3 = p.add_run()
                        run3.text = rest_text
                        run3.font.size = Pt(11)
                        run3.font.color.rgb = DARK_TEXT
                        run3.font.bold = False
                        run3.font.name = FONT_NAME
                else:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.space_before = Pt(3)
                    p.space_after = Pt(1)
                    run1 = p.add_run()
                    run1.text = "▸ "
                    run1.font.size = Pt(11)
                    run1.font.color.rgb = BLUE_PRIMARY
                    run1.font.bold = True
                    run1.font.name = FONT_NAME
                    run2 = p.add_run()
                    run2.text = rest_text
                    run2.font.size = Pt(11)
                    run2.font.color.rgb = DARK_TEXT
                    run2.font.name = FONT_NAME

            item_count = len(items)
            y_pos += Inches(0.25 * max(item_count, 1))

        # Ordered list
        elif el.name == "ol":
            lis = el.find_all("li", recursive=False)
            tf = add_textbox(slide, Inches(1.2), y_pos, CONTENT_WIDTH - Inches(0.4), Inches(0.3),
                             "", font_size=11, font_color=DARK_TEXT)
            tf.paragraphs[0].text = ""
            first = True
            for idx, li in enumerate(lis, 1):
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.space_before = Pt(2)
                p.space_after = Pt(1)
                run = p.add_run()
                run.text = f"{idx}. {li.get_text(strip=True)}"
                run.font.size = Pt(11)
                run.font.color.rgb = DARK_TEXT
                run.font.name = FONT_NAME
            y_pos += Inches(0.22 * max(len(lis), 1))

        # Paragraph text (standalone <p>)
        elif el.name == "p":
            text = el.get_text(strip=True)
            if text:
                tf = add_textbox(slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, Inches(0.4),
                                 text, font_size=11, font_color=DARK_TEXT)
                y_pos += Inches(0.35)

        # Highlight box
        elif "highlight-box" in classes:
            box_height = _estimate_box_height(el)
            # Background rectangle with left border
            add_colored_rect(slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, LIGHT_BLUE_BG)
            add_colored_rect(slide, LEFT_MARGIN, y_pos, Pt(5), box_height, BLUE_PRIMARY)

            inner_y = y_pos + Inches(0.1)
            h4 = el.find("h4")
            if h4:
                add_textbox(slide, Inches(1.0), inner_y, CONTENT_WIDTH - Inches(0.4), Inches(0.3),
                            h4.get_text(strip=True), font_size=13, font_color=BLUE_PRIMARY, bold=True)
                inner_y += Inches(0.3)

            # Content inside highlight box
            _render_box_content(slide, el, Inches(1.0), inner_y, CONTENT_WIDTH - Inches(0.4), DARK_TEXT)
            y_pos += box_height + Inches(0.1)

        # Warning box
        elif "warning-box" in classes:
            box_height = _estimate_box_height(el)
            add_colored_rect(slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, YELLOW_BG)
            add_colored_rect(slide, LEFT_MARGIN, y_pos, Pt(5), box_height, YELLOW_BORDER)

            inner_y = y_pos + Inches(0.1)
            strong = el.find("strong", recursive=False)
            if strong:
                add_textbox(slide, Inches(1.0), inner_y, CONTENT_WIDTH - Inches(0.4), Inches(0.3),
                            strong.get_text(strip=True), font_size=13, font_color=BROWN_TEXT, bold=True)
                inner_y += Inches(0.3)

            _render_box_content(slide, el, Inches(1.0), inner_y, CONTENT_WIDTH - Inches(0.4), BROWN_TEXT)
            y_pos += box_height + Inches(0.1)

        # Success box
        elif "success-box" in classes:
            box_height = _estimate_box_height(el)
            add_colored_rect(slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, GREEN_BG)
            add_colored_rect(slide, LEFT_MARGIN, y_pos, Pt(5), box_height, GREEN_BORDER)

            inner_y = y_pos + Inches(0.1)
            h4 = el.find("h4")
            if h4:
                add_textbox(slide, Inches(1.0), inner_y, CONTENT_WIDTH - Inches(0.4), Inches(0.3),
                            h4.get_text(strip=True), font_size=13, font_color=GREEN_DARK, bold=True)
                inner_y += Inches(0.3)

            _render_box_content(slide, el, Inches(1.0), inner_y, CONTENT_WIDTH - Inches(0.4), DARK_TEXT)
            y_pos += box_height + Inches(0.1)

        # Placeholder box
        elif "placeholder" in classes:
            box_height = _estimate_box_height(el, base=0.6)
            add_colored_rect(slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, PLACEHOLDER_BG, PLACEHOLDER_BORDER, 2)

            inner_y = y_pos + Inches(0.12)
            strong = el.find("strong")
            if strong:
                add_textbox(slide, LEFT_MARGIN, inner_y, CONTENT_WIDTH, Inches(0.3),
                            strong.get_text(strip=True), font_size=13, font_color=PLACEHOLDER_TEXT,
                            bold=True, alignment=PP_ALIGN.CENTER)
                inner_y += Inches(0.3)

            ems = el.find_all("em")
            if ems:
                tf = add_textbox(slide, Inches(1.2), inner_y, CONTENT_WIDTH - Inches(0.8), Inches(0.3),
                                 "", font_size=11, font_color=PLACEHOLDER_EM, italic=True,
                                 alignment=PP_ALIGN.CENTER)
                tf.paragraphs[0].text = ""
                first = True
                for em in ems:
                    text = em.get_text(strip=True)
                    if text:
                        if first:
                            tf.paragraphs[0].text = text
                            tf.paragraphs[0].font.size = Pt(11)
                            tf.paragraphs[0].font.color.rgb = PLACEHOLDER_EM
                            tf.paragraphs[0].font.italic = True
                            tf.paragraphs[0].font.name = FONT_NAME
                            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                            first = False
                        else:
                            add_paragraph(tf, text, font_size=11, font_color=PLACEHOLDER_EM,
                                          italic=True, alignment=PP_ALIGN.CENTER,
                                          space_before=Pt(2), space_after=Pt(1))

            y_pos += box_height + Inches(0.1)

        # Two-column layout
        elif "two-column" in classes:
            columns = el.select(".column-box")
            col_width = (CONTENT_WIDTH - Inches(0.3)) / 2

            max_col_h = Inches(0)
            for ci, col in enumerate(columns):
                col_left = LEFT_MARGIN + (col_width + Inches(0.3)) * ci
                col_h = _estimate_box_height(col, base=0.3)
                if col_h > max_col_h:
                    max_col_h = col_h

                # Column background
                add_colored_rect(slide, col_left, y_pos, col_width, col_h, GRAY_BG, GRAY_BORDER, 1)

                inner_y = y_pos + Inches(0.08)
                h4 = col.find("h4")
                if h4:
                    add_textbox(slide, col_left + Inches(0.15), inner_y, col_width - Inches(0.3), Inches(0.3),
                                h4.get_text(strip=True), font_size=12, font_color=BLUE_PRIMARY, bold=True)
                    inner_y += Inches(0.28)

                # Bullet list in column
                ul = col.find("ul", class_="bullet-list")
                if ul:
                    items = parse_bullet_items(ul)
                    _render_bullet_items_at(slide, items, col_left + Inches(0.15), inner_y,
                                            col_width - Inches(0.3), font_size=10)

                # Ordered list in column
                ol = col.find("ol")
                if ol:
                    lis = ol.find_all("li", recursive=False)
                    tf = add_textbox(slide, col_left + Inches(0.15), inner_y,
                                     col_width - Inches(0.3), Inches(0.3),
                                     "", font_size=10, font_color=DARK_TEXT)
                    tf.paragraphs[0].text = ""
                    first = True
                    for idx_li, li in enumerate(lis, 1):
                        if first:
                            tf.paragraphs[0].text = f"{idx_li}. {li.get_text(strip=True)}"
                            tf.paragraphs[0].font.size = Pt(10)
                            tf.paragraphs[0].font.color.rgb = DARK_TEXT
                            tf.paragraphs[0].font.name = FONT_NAME
                            first = False
                        else:
                            add_paragraph(tf, f"{idx_li}. {li.get_text(strip=True)}",
                                          font_size=10, font_color=DARK_TEXT,
                                          space_before=Pt(2), space_after=Pt(1))

            y_pos += max_col_h + Inches(0.12)

        # Comparison table
        elif el.name == "table" and "comparison-table" in classes:
            _render_table(slide, el, LEFT_MARGIN, y_pos, CONTENT_WIDTH)
            rows = el.find_all("tr")
            y_pos += Inches(0.3 * len(rows)) + Inches(0.1)

        i += 1

    return slide


def _estimate_box_height(el, base=0.4):
    """Estimate height of a box element based on content."""
    text = el.get_text(strip=True)
    lines = max(1, len(text) // 100 + text.count('\n') + 1)
    # Count bullet items
    lis = el.find_all("li")
    lines += len(lis)
    h4s = el.find_all("h4")
    lines += len(h4s)
    strongs = el.find_all("strong")
    # Count nested lists
    nested = el.find_all("ul")
    for n in nested:
        lines += len(n.find_all("li"))
    height = Inches(base + 0.18 * min(lines, 20))
    return min(height, Inches(3.5))


def _render_box_content(slide, box_el, left, y_pos, width, text_color):
    """Render bullet lists, paragraphs, etc. inside a colored box."""
    ul = box_el.find("ul", class_="bullet-list")
    if ul:
        items = parse_bullet_items(ul)
        _render_bullet_items_at(slide, items, left, y_pos, width, font_size=11, text_color=text_color)
        return

    # Paragraphs
    ps = box_el.find_all("p")
    if ps:
        tf = add_textbox(slide, left, y_pos, width, Inches(0.3), "", font_size=11, font_color=text_color)
        tf.paragraphs[0].text = ""
        first = True
        for p_el in ps:
            text = p_el.get_text(strip=True)
            if text:
                strong = p_el.find("strong")
                if first:
                    first = False
                    if strong:
                        bold_text = strong.get_text(strip=True)
                        rest = text[len(bold_text):]
                        run1 = tf.paragraphs[0].add_run()
                        run1.text = bold_text
                        run1.font.size = Pt(11)
                        run1.font.color.rgb = text_color
                        run1.font.bold = True
                        run1.font.name = FONT_NAME
                        if rest:
                            run2 = tf.paragraphs[0].add_run()
                            run2.text = rest
                            run2.font.size = Pt(11)
                            run2.font.color.rgb = text_color
                            run2.font.name = FONT_NAME
                    else:
                        tf.paragraphs[0].text = text
                        tf.paragraphs[0].font.size = Pt(11)
                        tf.paragraphs[0].font.color.rgb = text_color
                        tf.paragraphs[0].font.name = FONT_NAME
                else:
                    if strong:
                        bold_text = strong.get_text(strip=True)
                        rest = text[len(bold_text):]
                        p = tf.add_paragraph()
                        p.space_before = Pt(4)
                        run1 = p.add_run()
                        run1.text = bold_text
                        run1.font.size = Pt(11)
                        run1.font.color.rgb = text_color
                        run1.font.bold = True
                        run1.font.name = FONT_NAME
                        if rest:
                            run2 = p.add_run()
                            run2.text = rest
                            run2.font.size = Pt(11)
                            run2.font.color.rgb = text_color
                            run2.font.name = FONT_NAME
                    else:
                        add_paragraph(tf, text, font_size=11, font_color=text_color,
                                      space_before=Pt(4))


def _render_bullet_items_at(slide, items, left, y_pos, width, font_size=11, text_color=DARK_TEXT):
    """Render parsed bullet items at a given position."""
    tf = add_textbox(slide, left, y_pos, width, Inches(0.3), "", font_size=font_size, font_color=text_color)
    tf.paragraphs[0].text = ""
    first = True
    for bold_part, rest_text in items:
        if bold_part == "__nested__":
            p = tf.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(1)
            run = p.add_run()
            run.text = f"    → {rest_text}"
            run.font.size = Pt(font_size - 1)
            run.font.color.rgb = text_color
            run.font.name = FONT_NAME
        elif bold_part:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(1)
            run1 = p.add_run()
            run1.text = "▸ "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = BLUE_PRIMARY
            run1.font.bold = True
            run1.font.name = FONT_NAME
            run2 = p.add_run()
            run2.text = bold_part
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = text_color
            run2.font.bold = True
            run2.font.name = FONT_NAME
            if rest_text:
                run3 = p.add_run()
                run3.text = rest_text
                run3.font.size = Pt(font_size)
                run3.font.color.rgb = text_color
                run3.font.name = FONT_NAME
        else:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(1)
            run1 = p.add_run()
            run1.text = "▸ "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = BLUE_PRIMARY
            run1.font.bold = True
            run1.font.name = FONT_NAME
            run2 = p.add_run()
            run2.text = rest_text
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = text_color
            run2.font.name = FONT_NAME


def _render_table(slide, table_el, left, top, width):
    """Render an HTML table as a PowerPoint table shape."""
    rows = table_el.find_all("tr")
    if not rows:
        return

    header_row = rows[0]
    ths = header_row.find_all("th")
    num_cols = len(ths) if ths else len(rows[0].find_all("td"))
    num_rows = len(rows)

    if num_cols == 0 or num_rows == 0:
        return

    row_height = Inches(0.32)
    table_height = row_height * num_rows
    col_width = width // num_cols

    tbl_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, table_height)
    table = tbl_shape.table

    for row_idx, tr in enumerate(rows):
        cells = tr.find_all(["th", "td"])
        for col_idx, cell in enumerate(cells):
            if col_idx >= num_cols:
                break
            tbl_cell = table.cell(row_idx, col_idx)
            cell_text = cell.get_text(strip=True)

            # Check for checkmark cells
            has_check = "✓" in cell_text

            tbl_cell.text = ""
            p = tbl_cell.text_frame.paragraphs[0]

            # Bold parts
            strong = cell.find("strong")
            if strong and cell.name == "td":
                bold_text = strong.get_text(strip=True)
                rest = cell_text[len(bold_text):]
                run1 = p.add_run()
                run1.text = bold_text
                run1.font.bold = True
                run1.font.size = Pt(10)
                run1.font.name = FONT_NAME
                if rest:
                    run2 = p.add_run()
                    run2.text = rest
                    run2.font.size = Pt(10)
                    run2.font.name = FONT_NAME
            else:
                run = p.add_run()
                run.text = cell_text
                run.font.size = Pt(10)
                run.font.name = FONT_NAME

            # Header row styling
            if row_idx == 0:
                tbl_cell.fill.solid()
                tbl_cell.fill.fore_color.rgb = BLUE_PRIMARY
                for run in p.runs:
                    run.font.color.rgb = WHITE
                    run.font.bold = True
            else:
                # Check for green background (cells with ✓)
                style_attr = cell.get("style", "")
                if "dcfce7" in style_attr or has_check:
                    tbl_cell.fill.solid()
                    tbl_cell.fill.fore_color.rgb = GREEN_CELL_BG
                elif row_idx % 2 == 0:
                    tbl_cell.fill.solid()
                    tbl_cell.fill.fore_color.rgb = TABLE_EVEN_ROW
                for run in p.runs:
                    run.font.color.rgb = DARK_TEXT


# ────────────────────── Process all slides ──────────────────────

for idx, slide_html in enumerate(slides_html):
    classes = slide_html.get("class", [])
    if "title-slide" in classes:
        print(f"  Building title slide {idx + 1}...")
        build_title_slide(slide_html, prs)
    else:
        # Get slide number for logging
        num_el = slide_html.select_one(".slide-number")
        num = num_el.get_text(strip=True) if num_el else str(idx)
        print(f"  Building content slide {num}...")
        build_content_slide(slide_html, prs)

# ────────────────────── Save ──────────────────────

prs.save(str(OUTPUT_FILE))
print(f"\nSaved: {OUTPUT_FILE}")
print("Done!")
