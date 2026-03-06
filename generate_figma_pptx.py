"""
Generate a styled PPTX from FigmaToCodeAIAgentFramework2.html
Retains colors, fonts, styling, background colors, highlights, tables, and layout.
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image, ImageDraw
from bs4 import BeautifulSoup

BASE_DIR = Path(r"c:\Users\kumar.gn\PycharmProjects\Testproject")
HTML_FILE = BASE_DIR / "FigmaToCodeAIAgentFramework2.html"
OUTPUT_FILE = BASE_DIR / "FigmaToCodeAIAgentFramework2.pptx"
ASSETS_DIR = BASE_DIR / "_pptx_assets"
ASSETS_DIR.mkdir(exist_ok=True)

# Slide dimensions (widescreen 16:9)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors from HTML CSS
BLUE_PRIMARY = RGBColor(0x1E, 0x40, 0xAF)
BLUE_DARK = RGBColor(0x1E, 0x3A, 0x8A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK_TEXT = RGBColor(0x20, 0x20, 0x20)
LIGHT_BLUE_BG = RGBColor(0xDB, 0xEA, 0xFE)
LIGHT_BLUE_BG2 = RGBColor(0xBF, 0xDB, 0xFE)
YELLOW_BG = RGBColor(0xFE, 0xF0, 0x8A)
YELLOW_BORDER = RGBColor(0xCA, 0x8A, 0x04)
BROWN_TEXT = RGBColor(0x78, 0x35, 0x0F)
GREEN_BG = RGBColor(0xDC, 0xFC, 0xE7)
GREEN_BORDER = RGBColor(0x16, 0xA3, 0x4A)
GREEN_DARK = RGBColor(0x16, 0x65, 0x34)
PLACEHOLDER_BG = RGBColor(0xFE, 0xF3, 0xC7)
PLACEHOLDER_BORDER = RGBColor(0xD9, 0x77, 0x06)
PLACEHOLDER_TEXT = RGBColor(0x92, 0x40, 0x0E)
PLACEHOLDER_EM = RGBColor(0xB4, 0x53, 0x09)
GRAY_BG = RGBColor(0xF9, 0xFA, 0xFB)
GRAY_BORDER = RGBColor(0xE5, 0xE7, 0xEB)
TABLE_EVEN_ROW = RGBColor(0xF9, 0xFA, 0xFB)
SUBTITLE_COLOR = RGBColor(0xDB, 0xEA, 0xFE)
GREEN_CELL_BG = RGBColor(0xDC, 0xFC, 0xE7)

FONT_NAME = "Segoe UI"

# ────────────────────── Image helpers ──────────────────────


def make_gradient_image(width, height, color1, color2, filename):
    """Create a linear gradient image from color1 to color2 (diagonal 135deg)."""
    img = Image.new("RGB", (width, height))
    for y in range(height):
        for x in range(width):
            t = (x / width + y / height) / 2.0
            r = int(color1[0] * (1 - t) + color2[0] * t)
            g = int(color1[1] * (1 - t) + color2[1] * t)
            b = int(color1[2] * (1 - t) + color2[2] * t)
            img.putpixel((x, y), (r, g, b))
    path = ASSETS_DIR / filename
    img.save(str(path))
    return str(path)


def make_solid_image(width, height, color, filename):
    img = Image.new("RGB", (width, height), color)
    path = ASSETS_DIR / filename
    img.save(str(path))
    return str(path)


# Generate background images
print("Generating background images...")
title_bg_path = make_gradient_image(
    400, 225, (0x1E, 0x40, 0xAF), (0x1E, 0x3A, 0x8A), "figma_title_bg.png"
)
highlight_bg_path = make_gradient_image(
    400, 60, (0xDB, 0xEA, 0xFE), (0xBF, 0xDB, 0xFE), "figma_highlight_bg.png"
)
warning_bg_path = make_gradient_image(
    400, 60, (0xFE, 0xF0, 0x8A), (0xFD, 0xE0, 0x47), "figma_warning_bg.png"
)
success_bg_path = make_gradient_image(
    400, 60, (0xDC, 0xFC, 0xE7), (0xBB, 0xF7, 0xD0), "figma_success_bg.png"
)
white_bg_path = make_solid_image(400, 225, (0xFF, 0xFF, 0xFF), "figma_white_bg.png")

# ────────────────────── PPTX helpers ──────────────────────


def add_bg_picture(slide, image_path):
    """Add a background image covering the full slide."""
    pic = slide.shapes.add_picture(image_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text="",
    font_size=14,
    font_color=DARK_TEXT,
    bold=False,
    italic=False,
    alignment=PP_ALIGN.LEFT,
    font_name=FONT_NAME,
):
    """Add a text box and return the text frame for further manipulation."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(
    tf,
    text="",
    font_size=14,
    font_color=DARK_TEXT,
    bold=False,
    italic=False,
    alignment=PP_ALIGN.LEFT,
    space_before=Pt(4),
    space_after=Pt(2),
    font_name=FONT_NAME,
    level=0,
):
    """Add a paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = bold
    p.font.italic = italic
    p.font.name = font_name
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    return p


def add_colored_rect(
    slide, left, top, width, height, fill_color, border_color=None, border_width=None
):
    """Add a colored rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(border_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_circle(
    slide, left, top, size, fill_color, text="", font_size=24, font_color=WHITE
):
    """Add a circle with centered text."""
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = font_color
    p.font.bold = True
    p.font.name = FONT_NAME
    p.alignment = PP_ALIGN.CENTER
    from pptx.oxml.ns import qn

    txBody = shape.text_frame._txBody
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.set("anchor", "ctr")
    return shape


# ────────────────────── Parse HTML ──────────────────────

print("Parsing HTML...")
with open(HTML_FILE, "r", encoding="utf-8") as f:
    soup = BeautifulSoup(f.read(), "html.parser")

slides_html = soup.select(".slide")

# ────────────────────── HTML content parser helpers ──────────────────────


def parse_bullet_items(ul_el):
    """Parse bullet list items, returning list of (bold_part, rest_text) tuples.
    bold_part is None for plain items, '__nested__' for sub-items."""
    items = []
    if not ul_el:
        return items
    for li in ul_el.find_all("li", recursive=False):
        strong = li.find("strong", recursive=False)
        nested_ul = li.find("ul", class_="bullet-list")

        if strong:
            bold_text = strong.get_text(strip=True)
            full_text = li.get_text(strip=True)
            rest = full_text[len(bold_text) :] if full_text.startswith(bold_text) else full_text.replace(bold_text, "", 1)
            items.append((bold_text, rest))
        else:
            # Get main text (before any nested list)
            text_parts = []
            for child in li.children:
                if hasattr(child, "name") and child.name == "ul":
                    break
                if hasattr(child, "get_text"):
                    t = child.get_text(strip=True)
                    if t:
                        text_parts.append(t)
                elif isinstance(child, str) and child.strip():
                    text_parts.append(child.strip())
            main_text = " ".join(text_parts) if text_parts else li.get_text(strip=True)
            items.append((None, main_text))

        # Process nested bullet items
        if nested_ul:
            for nested_li in nested_ul.find_all("li", recursive=False):
                nested_strong = nested_li.find("strong")
                if nested_strong:
                    nb = nested_strong.get_text(strip=True)
                    nfull = nested_li.get_text(strip=True)
                    nrest = nfull[len(nb):] if nfull.startswith(nb) else nfull.replace(nb, "", 1)
                    items.append(("__nested_bold__", (nb, nrest)))
                else:
                    items.append(("__nested__", nested_li.get_text(strip=True)))
    return items


# ────────────────────── Height estimation helpers ──────────────────────

def _estimate_lines(text, width_inches, font_pt):
    """Estimate number of rendered lines for text given width and font size.
    Uses conservative Segoe UI character width estimates."""
    # Segoe UI average chars per inch:  ~13 at 9pt, ~12 at 10pt, ~14.5 at 8pt
    chars_per_inch = 115.0 / font_pt
    chars_per_line = max(10, int(width_inches * chars_per_inch))
    text_len = len(text)
    return max(1, -(-text_len // chars_per_line))  # ceiling division


def _line_height(font_pt):
    """Line height in inches including spacing for a given font pt size."""
    return font_pt * 1.55 / 72.0


def _bullet_list_height(items, width_inches, font_pt=9):
    """Compute total height of a parsed bullet item list."""
    lh = _line_height(font_pt)
    nested_lh = _line_height(font_pt - 1)
    total = 0.0
    for bold_part, rest_text in items:
        if bold_part in ("__nested__", "__nested_bold__"):
            if bold_part == "__nested_bold__":
                nb, nrest = rest_text
                text = f"    -> {nb}{nrest}"
            else:
                text = f"    -> {rest_text}"
            lines = _estimate_lines(text, width_inches, font_pt - 1)
            total += lines * nested_lh + 0.02
        else:
            full = f"x {bold_part or ''}{rest_text}"  # "▸ " prefix
            lines = _estimate_lines(full, width_inches, font_pt)
            total += lines * lh + 0.02
    return total


def _estimate_box_height(el, width_inches=None, base=0.30):
    """Estimate height of a box element based on actual text content."""
    if width_inches is None:
        width_inches = 11.0  # CONTENT_WIDTH - margins, approx
    total = base  # padding top+bottom
    # h4 heading
    for h4 in el.find_all("h4"):
        total += 0.24
    # paragraphs
    for p in el.find_all("p"):
        text = p.get_text(strip=True)
        if text:
            lines = _estimate_lines(text, width_inches, 9)
            total += lines * _line_height(9) + 0.02
    # bullet items
    uls = el.find_all("ul", class_="bullet-list")
    for ul in uls:
        items = parse_bullet_items(ul)
        total += _bullet_list_height(items, width_inches, 9)
    return Inches(min(total, 3.5))


def _render_box_content(slide, box_el, left, y_pos, width, text_color):
    """Render bullet lists, paragraphs, etc. inside a colored box."""
    ul = box_el.find("ul", class_="bullet-list")
    if ul:
        items = parse_bullet_items(ul)
        _render_bullet_items_at(
            slide, items, left, y_pos, width, font_size=9, text_color=text_color
        )
        return

    # Paragraphs
    ps = box_el.find_all("p")
    if ps:
        tf = add_textbox(
            slide, left, y_pos, width, Inches(0.3), "", font_size=9, font_color=text_color
        )
        tf.paragraphs[0].text = ""
        first = True
        for p_el in ps:
            text = p_el.get_text(strip=True)
            if text:
                strong = p_el.find("strong")
                if first:
                    first = False
                    if strong:
                        bold_text = strong.get_text(strip=True)
                        rest = text[len(bold_text) :]
                        run1 = tf.paragraphs[0].add_run()
                        run1.text = bold_text
                        run1.font.size = Pt(9)
                        run1.font.color.rgb = text_color
                        run1.font.bold = True
                        run1.font.name = FONT_NAME
                        if rest:
                            run2 = tf.paragraphs[0].add_run()
                            run2.text = rest
                            run2.font.size = Pt(9)
                            run2.font.color.rgb = text_color
                            run2.font.name = FONT_NAME
                    else:
                        tf.paragraphs[0].text = text
                        tf.paragraphs[0].font.size = Pt(9)
                        tf.paragraphs[0].font.color.rgb = text_color
                        tf.paragraphs[0].font.name = FONT_NAME
                else:
                    if strong:
                        bold_text = strong.get_text(strip=True)
                        rest = text[len(bold_text) :]
                        p = tf.add_paragraph()
                        p.space_before = Pt(2)
                        run1 = p.add_run()
                        run1.text = bold_text
                        run1.font.size = Pt(9)
                        run1.font.color.rgb = text_color
                        run1.font.bold = True
                        run1.font.name = FONT_NAME
                        if rest:
                            run2 = p.add_run()
                            run2.text = rest
                            run2.font.size = Pt(9)
                            run2.font.color.rgb = text_color
                            run2.font.name = FONT_NAME
                    else:
                        add_paragraph(
                            tf,
                            text,
                            font_size=9,
                            font_color=text_color,
                            space_before=Pt(2),
                        )


def _render_bullet_items_at(
    slide, items, left, y_pos, width, font_size=10, text_color=DARK_TEXT
):
    """Render parsed bullet items at a given position."""
    width_in = width / 914400  # EMU to inches
    box_h = _bullet_list_height(items, width_in, font_size)
    tf = add_textbox(
        slide, left, y_pos, width, Inches(max(0.3, box_h)), "", font_size=font_size, font_color=text_color
    )
    tf.paragraphs[0].text = ""
    first = True
    for bold_part, rest_text in items:
        if bold_part == "__nested__":
            p = tf.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(1)
            run = p.add_run()
            run.text = f"    → {rest_text}"
            run.font.size = Pt(font_size - 1)
            run.font.color.rgb = text_color
            run.font.name = FONT_NAME
        elif bold_part == "__nested_bold__":
            nb, nrest = rest_text
            p = tf.add_paragraph()
            p.space_before = Pt(1)
            p.space_after = Pt(1)
            run0 = p.add_run()
            run0.text = "    → "
            run0.font.size = Pt(font_size - 1)
            run0.font.color.rgb = text_color
            run0.font.name = FONT_NAME
            run1 = p.add_run()
            run1.text = nb
            run1.font.size = Pt(font_size - 1)
            run1.font.color.rgb = text_color
            run1.font.bold = True
            run1.font.name = FONT_NAME
            if nrest:
                run2 = p.add_run()
                run2.text = nrest
                run2.font.size = Pt(font_size - 1)
                run2.font.color.rgb = text_color
                run2.font.name = FONT_NAME
        elif bold_part:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(1)
            run1 = p.add_run()
            run1.text = "▸ "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = BLUE_PRIMARY
            run1.font.bold = True
            run1.font.name = FONT_NAME
            run2 = p.add_run()
            run2.text = bold_part
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = text_color
            run2.font.bold = True
            run2.font.name = FONT_NAME
            if rest_text:
                run3 = p.add_run()
                run3.text = rest_text
                run3.font.size = Pt(font_size)
                run3.font.color.rgb = text_color
                run3.font.name = FONT_NAME
        else:
            if first:
                p = tf.paragraphs[0]
                first = False
            else:
                p = tf.add_paragraph()
            p.space_before = Pt(2)
            p.space_after = Pt(1)
            run1 = p.add_run()
            run1.text = "▸ "
            run1.font.size = Pt(font_size)
            run1.font.color.rgb = BLUE_PRIMARY
            run1.font.bold = True
            run1.font.name = FONT_NAME
            run2 = p.add_run()
            run2.text = rest_text
            run2.font.size = Pt(font_size)
            run2.font.color.rgb = text_color
            run2.font.name = FONT_NAME


def _render_table(slide, table_el, left, top, width):
    """Render an HTML table as a PowerPoint table shape."""
    rows = table_el.find_all("tr")
    if not rows:
        return

    header_row = rows[0]
    ths = header_row.find_all("th")
    num_cols = len(ths) if ths else len(rows[0].find_all("td"))
    num_rows = len(rows)

    if num_cols == 0 or num_rows == 0:
        return

    row_height = Inches(0.27)
    table_height = row_height * num_rows
    col_width = width // num_cols

    tbl_shape = slide.shapes.add_table(
        num_rows, num_cols, left, top, width, table_height
    )
    table = tbl_shape.table

    for row_idx, tr in enumerate(rows):
        cells = tr.find_all(["th", "td"])
        for col_idx, cell in enumerate(cells):
            if col_idx >= num_cols:
                break
            tbl_cell = table.cell(row_idx, col_idx)
            cell_text = cell.get_text(strip=True)
            has_check = "✓" in cell_text

            tbl_cell.text = ""
            p = tbl_cell.text_frame.paragraphs[0]

            strong = cell.find("strong")
            if strong and cell.name == "td":
                bold_text = strong.get_text(strip=True)
                rest = cell_text[len(bold_text) :]
                run1 = p.add_run()
                run1.text = bold_text
                run1.font.bold = True
                run1.font.size = Pt(8)
                run1.font.name = FONT_NAME
                if rest:
                    run2 = p.add_run()
                    run2.text = rest
                    run2.font.size = Pt(8)
                    run2.font.name = FONT_NAME
            else:
                run = p.add_run()
                run.text = cell_text
                run.font.size = Pt(8)
                run.font.name = FONT_NAME

            if row_idx == 0:
                tbl_cell.fill.solid()
                tbl_cell.fill.fore_color.rgb = BLUE_PRIMARY
                for run in p.runs:
                    run.font.color.rgb = WHITE
                    run.font.bold = True
            else:
                style_attr = cell.get("style", "")
                if "dcfce7" in style_attr or has_check:
                    tbl_cell.fill.solid()
                    tbl_cell.fill.fore_color.rgb = GREEN_CELL_BG
                elif row_idx % 2 == 0:
                    tbl_cell.fill.solid()
                    tbl_cell.fill.fore_color.rgb = TABLE_EVEN_ROW
                for run in p.runs:
                    run.font.color.rgb = DARK_TEXT


# ────────────────────── Build PPTX ──────────────────────

print("Building PPTX...")
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]

# Margins
LEFT_MARGIN = Inches(0.8)
TOP_MARGIN = Inches(0.3)
CONTENT_WIDTH = SLIDE_W - Inches(1.6)
HEADER_HEIGHT = Inches(1.1)
BODY_TOP = Inches(1.5)


def build_title_slide(slide_html, prs):
    """Build a title slide (blue gradient background, centered white text)."""
    slide = prs.slides.add_slide(blank_layout)
    add_bg_picture(slide, title_bg_path)

    title_el = slide_html.select_one(".slide-title")
    paragraphs = slide_html.select("p")

    title_text = title_el.get_text(strip=True) if title_el else ""
    subtitles = [p.get_text(strip=True) for p in paragraphs if p.get_text(strip=True)]

    # Title
    tf = add_textbox(
        slide,
        LEFT_MARGIN,
        Inches(2.2),
        CONTENT_WIDTH,
        Inches(1.5),
        title_text,
        font_size=44,
        font_color=WHITE,
        bold=True,
        alignment=PP_ALIGN.CENTER,
    )

    # Subtitles
    for st in subtitles:
        add_paragraph(
            tf,
            st,
            font_size=20,
            font_color=SUBTITLE_COLOR,
            alignment=PP_ALIGN.CENTER,
            space_before=Pt(12),
            space_after=Pt(6),
        )

    return slide


def _add_slide_header(slide, slide_num, slide_title):
    """Add the standard header (circle, title, underline) to a slide."""
    add_colored_rect(slide, LEFT_MARGIN, Inches(1.05), CONTENT_WIDTH, Pt(3), BLUE_PRIMARY)
    circle_size = Inches(0.55)
    add_circle(
        slide, LEFT_MARGIN, Inches(0.3), circle_size, BLUE_PRIMARY,
        text=slide_num, font_size=18, font_color=WHITE,
    )
    add_textbox(
        slide, Inches(1.55), Inches(0.25), CONTENT_WIDTH - Inches(0.75), Inches(0.75),
        slide_title, font_size=20, font_color=BLUE_PRIMARY, bold=True,
    )


def build_content_slide(slide_html, prs):
    """Build content slide(s) with auto-continuation if content overflows."""
    slide = prs.slides.add_slide(blank_layout)
    add_bg_picture(slide, white_bg_path)

    number_el = slide_html.select_one(".slide-number")
    title_el = slide_html.select_one(".slide-title")
    slide_num = number_el.get_text(strip=True) if number_el else ""
    slide_title = title_el.get_text(strip=True) if title_el else ""

    _add_slide_header(slide, slide_num, slide_title)

    # Process content blocks
    content_div = slide_html.select_one(".content")
    if not content_div:
        return slide

    y_pos = Inches(1.20)
    MAX_Y = Inches(7.05)  # Don't place content below this
    CONT_START_Y = Inches(1.20)  # y_pos for continuation slides

    def _need_continuation(needed_inches):
        """Check if we need a continuation slide for the next element."""
        nonlocal slide, y_pos
        if y_pos + Inches(needed_inches) > MAX_Y:
            # Create continuation slide
            slide = prs.slides.add_slide(blank_layout)
            add_bg_picture(slide, white_bg_path)
            _add_slide_header(slide, slide_num, slide_title + " (cont.)")
            y_pos = CONT_START_Y

    children = list(content_div.children)
    i = 0
    while i < len(children):
        el = children[i]
        if not hasattr(el, "name") or el.name is None:
            i += 1
            continue

        classes = el.get("class", [])

        # Section Heading
        if "section-heading" in classes:
            _need_continuation(0.27)
            heading_text = el.get_text(strip=True)
            tf = add_textbox(
                slide,
                LEFT_MARGIN,
                y_pos,
                CONTENT_WIDTH,
                Inches(0.26),
                heading_text,
                font_size=12,
                font_color=BLUE_PRIMARY,
                bold=True,
            )
            y_pos += Inches(0.27)

        # Bullet list (standalone)
        elif el.name == "ul" and "bullet-list" in classes:
            items = parse_bullet_items(el)
            avail_w_in = (CONTENT_WIDTH - Inches(0.3)) / 914400
            bullet_h = _bullet_list_height(items, avail_w_in, 9)
            _need_continuation(bullet_h)
            tf = add_textbox(
                slide,
                Inches(1.1),
                y_pos,
                CONTENT_WIDTH - Inches(0.3),
                Inches(max(0.3, bullet_h)),
                "",
                font_size=9,
                font_color=DARK_TEXT,
            )
            tf.paragraphs[0].text = ""
            first = True
            for bold_part, rest_text in items:
                if bold_part == "__nested__":
                    p = tf.add_paragraph()
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    run = p.add_run()
                    run.text = f"    → {rest_text}"
                    run.font.size = Pt(8)
                    run.font.color.rgb = DARK_TEXT
                    run.font.name = FONT_NAME
                elif bold_part == "__nested_bold__":
                    nb, nrest = rest_text
                    p = tf.add_paragraph()
                    p.space_before = Pt(0)
                    p.space_after = Pt(0)
                    run0 = p.add_run()
                    run0.text = "    → "
                    run0.font.size = Pt(8)
                    run0.font.color.rgb = DARK_TEXT
                    run0.font.name = FONT_NAME
                    run1 = p.add_run()
                    run1.text = nb
                    run1.font.size = Pt(8)
                    run1.font.color.rgb = DARK_TEXT
                    run1.font.bold = True
                    run1.font.name = FONT_NAME
                    if nrest:
                        run2 = p.add_run()
                        run2.text = nrest
                        run2.font.size = Pt(8)
                        run2.font.color.rgb = DARK_TEXT
                        run2.font.name = FONT_NAME
                elif bold_part:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.space_before = Pt(1)
                    p.space_after = Pt(0)
                    run1 = p.add_run()
                    run1.text = "▸ "
                    run1.font.size = Pt(9)
                    run1.font.color.rgb = BLUE_PRIMARY
                    run1.font.bold = True
                    run1.font.name = FONT_NAME
                    run2 = p.add_run()
                    run2.text = bold_part
                    run2.font.size = Pt(9)
                    run2.font.color.rgb = DARK_TEXT
                    run2.font.bold = True
                    run2.font.name = FONT_NAME
                    if rest_text:
                        run3 = p.add_run()
                        run3.text = rest_text
                        run3.font.size = Pt(9)
                        run3.font.color.rgb = DARK_TEXT
                        run3.font.bold = False
                        run3.font.name = FONT_NAME
                else:
                    if first:
                        p = tf.paragraphs[0]
                        first = False
                    else:
                        p = tf.add_paragraph()
                    p.space_before = Pt(1)
                    p.space_after = Pt(0)
                    run1 = p.add_run()
                    run1.text = "▸ "
                    run1.font.size = Pt(9)
                    run1.font.color.rgb = BLUE_PRIMARY
                    run1.font.bold = True
                    run1.font.name = FONT_NAME
                    run2 = p.add_run()
                    run2.text = rest_text
                    run2.font.size = Pt(9)
                    run2.font.color.rgb = DARK_TEXT
                    run2.font.name = FONT_NAME

            # Compute actual height based on text wrapping
            avail_w = (CONTENT_WIDTH - Inches(0.3)) / 914400  # to inches
            y_pos += Inches(_bullet_list_height(items, avail_w, 9))

        # Ordered list
        elif el.name == "ol":
            lis = el.find_all("li", recursive=False)
            ol_avail_w = (CONTENT_WIDTH - Inches(0.3)) / 914400
            ol_h_est = 0.0
            for li in lis:
                t = li.get_text(strip=True)
                ol_h_est += _estimate_lines(f"1. {t}", ol_avail_w, 9) * _line_height(9) + 0.02
            _need_continuation(ol_h_est)
            tf = add_textbox(
                slide,
                Inches(1.1),
                y_pos,
                CONTENT_WIDTH - Inches(0.3),
                Inches(max(0.3, ol_h_est)),
                "",
                font_size=9,
                font_color=DARK_TEXT,
            )
            tf.paragraphs[0].text = ""
            first = True
            for idx_li, li in enumerate(lis, 1):
                val = li.get("value")
                num = int(val) if val else idx_li
                if first:
                    p = tf.paragraphs[0]
                    first = False
                else:
                    p = tf.add_paragraph()
                p.space_before = Pt(2)
                p.space_after = Pt(1)
                strong = li.find("strong")
                if strong:
                    run0 = p.add_run()
                    run0.text = f"{num}. "
                    run0.font.size = Pt(9)
                    run0.font.color.rgb = DARK_TEXT
                    run0.font.name = FONT_NAME
                    run1 = p.add_run()
                    run1.text = strong.get_text(strip=True)
                    run1.font.size = Pt(9)
                    run1.font.color.rgb = DARK_TEXT
                    run1.font.bold = True
                    run1.font.name = FONT_NAME
                    rest = li.get_text(strip=True)[len(strong.get_text(strip=True)):]
                    if rest:
                        run2 = p.add_run()
                        run2.text = rest
                        run2.font.size = Pt(9)
                        run2.font.color.rgb = DARK_TEXT
                        run2.font.name = FONT_NAME
                else:
                    run = p.add_run()
                    run.text = f"{num}. {li.get_text(strip=True)}"
                    run.font.size = Pt(9)
                    run.font.color.rgb = DARK_TEXT
                    run.font.name = FONT_NAME
            # Compute actual height based on text
            avail_w = (CONTENT_WIDTH - Inches(0.3)) / 914400
            total_ol_h = 0.0
            ol_lh = _line_height(9)
            for li in lis:
                t = li.get_text(strip=True)
                lines = _estimate_lines(f"1. {t}", avail_w, 9)
                total_ol_h += lines * ol_lh + 0.02
            y_pos += Inches(total_ol_h)

        # Paragraph text (standalone <p>)
        elif el.name == "p":
            text = el.get_text(strip=True)
            if text:
                p_w = CONTENT_WIDTH / 914400
                p_lines = _estimate_lines(text, p_w, 9)
                p_h = max(0.25, p_lines * _line_height(9) + 0.04)
                _need_continuation(p_h)
                tf = add_textbox(
                    slide,
                    LEFT_MARGIN,
                    y_pos,
                    CONTENT_WIDTH,
                    Inches(p_h),
                    text,
                    font_size=9,
                    font_color=DARK_TEXT,
                )
                y_pos += Inches(p_h)

        # Highlight box
        elif "highlight-box" in classes:
            box_w = (CONTENT_WIDTH - Inches(0.4)) / 914400
            box_height = _estimate_box_height(el, width_inches=box_w)
            _need_continuation(box_height / 914400 + 0.06)
            add_colored_rect(
                slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, LIGHT_BLUE_BG
            )
            add_colored_rect(
                slide, LEFT_MARGIN, y_pos, Pt(5), box_height, BLUE_PRIMARY
            )

            inner_y = y_pos + Inches(0.06)
            h4 = el.find("h4")
            if h4:
                add_textbox(
                    slide,
                    Inches(1.0),
                    inner_y,
                    CONTENT_WIDTH - Inches(0.4),
                    Inches(0.22),
                    h4.get_text(strip=True),
                    font_size=11,
                    font_color=BLUE_PRIMARY,
                    bold=True,
                )
                inner_y += Inches(0.22)

            _render_box_content(
                slide,
                el,
                Inches(1.0),
                inner_y,
                CONTENT_WIDTH - Inches(0.4),
                DARK_TEXT,
            )
            y_pos += box_height + Inches(0.06)

        # Warning box
        elif "warning-box" in classes:
            box_w = (CONTENT_WIDTH - Inches(0.4)) / 914400
            box_height = _estimate_box_height(el, width_inches=box_w)
            _need_continuation(box_height / 914400 + 0.06)
            add_colored_rect(
                slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, YELLOW_BG
            )
            add_colored_rect(
                slide, LEFT_MARGIN, y_pos, Pt(5), box_height, YELLOW_BORDER
            )

            inner_y = y_pos + Inches(0.06)
            strong = el.find("strong", recursive=False)
            if strong:
                add_textbox(
                    slide,
                    Inches(1.0),
                    inner_y,
                    CONTENT_WIDTH - Inches(0.4),
                    Inches(0.22),
                    strong.get_text(strip=True),
                    font_size=11,
                    font_color=BROWN_TEXT,
                    bold=True,
                )
                inner_y += Inches(0.22)

            _render_box_content(
                slide,
                el,
                Inches(1.0),
                inner_y,
                CONTENT_WIDTH - Inches(0.4),
                BROWN_TEXT,
            )
            y_pos += box_height + Inches(0.06)

        # Success box
        elif "success-box" in classes:
            box_w = (CONTENT_WIDTH - Inches(0.4)) / 914400
            box_height = _estimate_box_height(el, width_inches=box_w)
            _need_continuation(box_height / 914400 + 0.06)
            add_colored_rect(
                slide, LEFT_MARGIN, y_pos, CONTENT_WIDTH, box_height, GREEN_BG
            )
            add_colored_rect(
                slide, LEFT_MARGIN, y_pos, Pt(5), box_height, GREEN_BORDER
            )

            inner_y = y_pos + Inches(0.06)
            h4 = el.find("h4")
            if h4:
                add_textbox(
                    slide,
                    Inches(1.0),
                    inner_y,
                    CONTENT_WIDTH - Inches(0.4),
                    Inches(0.22),
                    h4.get_text(strip=True),
                    font_size=11,
                    font_color=GREEN_DARK,
                    bold=True,
                )
                inner_y += Inches(0.22)

            _render_box_content(
                slide,
                el,
                Inches(1.0),
                inner_y,
                CONTENT_WIDTH - Inches(0.4),
                DARK_TEXT,
            )
            y_pos += box_height + Inches(0.06)

        # Placeholder box
        elif "placeholder" in classes:
            box_w = CONTENT_WIDTH / 914400
            box_height = _estimate_box_height(el, width_inches=box_w, base=0.4)
            _need_continuation(box_height / 914400 + 0.06)
            add_colored_rect(
                slide,
                LEFT_MARGIN,
                y_pos,
                CONTENT_WIDTH,
                box_height,
                PLACEHOLDER_BG,
                PLACEHOLDER_BORDER,
                2,
            )

            inner_y = y_pos + Inches(0.06)
            strong = el.find("strong")
            if strong:
                add_textbox(
                    slide,
                    LEFT_MARGIN,
                    inner_y,
                    CONTENT_WIDTH,
                    Inches(0.22),
                    strong.get_text(strip=True),
                    font_size=11,
                    font_color=PLACEHOLDER_TEXT,
                    bold=True,
                    alignment=PP_ALIGN.CENTER,
                )
                inner_y += Inches(0.22)

            ems = el.find_all("em")
            if ems:
                tf = add_textbox(
                    slide,
                    Inches(1.2),
                    inner_y,
                    CONTENT_WIDTH - Inches(0.8),
                    Inches(0.3),
                    "",
                    font_size=9,
                    font_color=PLACEHOLDER_EM,
                    italic=True,
                    alignment=PP_ALIGN.CENTER,
                )
                tf.paragraphs[0].text = ""
                first = True
                for em in ems:
                    text = em.get_text(strip=True)
                    if text:
                        if first:
                            tf.paragraphs[0].text = text
                            tf.paragraphs[0].font.size = Pt(9)
                            tf.paragraphs[0].font.color.rgb = PLACEHOLDER_EM
                            tf.paragraphs[0].font.italic = True
                            tf.paragraphs[0].font.name = FONT_NAME
                            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
                            first = False
                        else:
                            add_paragraph(
                                tf,
                                text,
                                font_size=9,
                                font_color=PLACEHOLDER_EM,
                                italic=True,
                                alignment=PP_ALIGN.CENTER,
                                space_before=Pt(1),
                                space_after=Pt(0),
                            )

            y_pos += box_height + Inches(0.06)

        # Two-column layout
        elif "two-column" in classes:
            columns = el.select(".column-box")
            col_width = (CONTENT_WIDTH - Inches(0.3)) / 2
            # Pre-compute max column height for continuation check
            pre_max_h = Inches(0)
            for col in columns:
                col_inner_w_pre = (col_width - Inches(0.24)) / 914400
                ch = _estimate_box_height(col, width_inches=col_inner_w_pre, base=0.25)
                if ch > pre_max_h:
                    pre_max_h = ch
            _need_continuation(pre_max_h / 914400 + 0.08)

            max_col_h = Inches(0)
            for ci, col in enumerate(columns):
                col_left = LEFT_MARGIN + (col_width + Inches(0.3)) * ci
                col_inner_w = (col_width - Inches(0.24)) / 914400
                col_h = _estimate_box_height(col, width_inches=col_inner_w, base=0.25)
                if col_h > max_col_h:
                    max_col_h = col_h

                add_colored_rect(
                    slide, col_left, y_pos, col_width, col_h, GRAY_BG, GRAY_BORDER, 1
                )

                inner_y = y_pos + Inches(0.05)
                h4 = col.find("h4")
                if h4:
                    add_textbox(
                        slide,
                        col_left + Inches(0.12),
                        inner_y,
                        col_width - Inches(0.24),
                        Inches(0.22),
                        h4.get_text(strip=True),
                        font_size=10,
                        font_color=BLUE_PRIMARY,
                        bold=True,
                    )
                    inner_y += Inches(0.21)

                # Paragraph text in column
                p_el = col.find("p")
                if p_el and not col.find("ul") and not col.find("ol"):
                    add_textbox(
                        slide,
                        col_left + Inches(0.12),
                        inner_y,
                        col_width - Inches(0.24),
                        Inches(0.45),
                        p_el.get_text(strip=True),
                        font_size=8,
                        font_color=DARK_TEXT,
                    )

                # Bullet list in column
                ul = col.find("ul", class_="bullet-list")
                if ul:
                    items = parse_bullet_items(ul)
                    _render_bullet_items_at(
                        slide,
                        items,
                        col_left + Inches(0.12),
                        inner_y,
                        col_width - Inches(0.24),
                        font_size=8,
                    )

                # Ordered list in column
                ol = col.find("ol")
                if ol:
                    lis = ol.find_all("li", recursive=False)
                    tf = add_textbox(
                        slide,
                        col_left + Inches(0.12),
                        inner_y,
                        col_width - Inches(0.24),
                        Inches(0.3),
                        "",
                        font_size=8,
                        font_color=DARK_TEXT,
                    )
                    tf.paragraphs[0].text = ""
                    first = True
                    for idx_li, li in enumerate(lis, 1):
                        val = li.get("value")
                        num = int(val) if val else idx_li
                        if first:
                            p = tf.paragraphs[0]
                            first = False
                        else:
                            p = tf.add_paragraph()
                        p.space_before = Pt(1)
                        p.space_after = Pt(0)
                        strong_li = li.find("strong")
                        if strong_li:
                            run0 = p.add_run()
                            run0.text = f"{num}. "
                            run0.font.size = Pt(8)
                            run0.font.color.rgb = DARK_TEXT
                            run0.font.name = FONT_NAME
                            run1 = p.add_run()
                            run1.text = strong_li.get_text(strip=True)
                            run1.font.size = Pt(8)
                            run1.font.color.rgb = DARK_TEXT
                            run1.font.bold = True
                            run1.font.name = FONT_NAME
                            rest = li.get_text(strip=True)[len(strong_li.get_text(strip=True)):]
                            if rest:
                                run2 = p.add_run()
                                run2.text = rest
                                run2.font.size = Pt(8)
                                run2.font.color.rgb = DARK_TEXT
                                run2.font.name = FONT_NAME
                        else:
                            run = p.add_run()
                            run.text = f"{num}. {li.get_text(strip=True)}"
                            run.font.size = Pt(8)
                            run.font.color.rgb = DARK_TEXT
                            run.font.name = FONT_NAME

            y_pos += max_col_h + Inches(0.08)

        # Comparison table
        elif el.name == "table" and "comparison-table" in classes:
            rows_list = el.find_all("tr")
            num_cols_est = max(1, len(rows_list[0].find_all(["th", "td"]))) if rows_list else 1
            col_w_in = (CONTENT_WIDTH / 914400) / num_cols_est
            tbl_h = 0.0
            for tr in rows_list:
                max_cell_lines = 1
                for cell in tr.find_all(["th", "td"]):
                    cl = _estimate_lines(cell.get_text(strip=True), col_w_in, 8)
                    if cl > max_cell_lines:
                        max_cell_lines = cl
                tbl_h += max_cell_lines * _line_height(8) + 0.06
            _need_continuation(tbl_h + 0.06)
            _render_table(slide, el, LEFT_MARGIN, y_pos, CONTENT_WIDTH)
            y_pos += Inches(tbl_h) + Inches(0.06)

        i += 1

    return slide


# ────────────────────── Process all slides ──────────────────────

for idx, slide_html in enumerate(slides_html):
    classes = slide_html.get("class", [])
    if "title-slide" in classes:
        print(f"  Building title slide {idx + 1}...")
        build_title_slide(slide_html, prs)
    else:
        num_el = slide_html.select_one(".slide-number")
        num = num_el.get_text(strip=True) if num_el else str(idx)
        print(f"  Building content slide {num}...")
        build_content_slide(slide_html, prs)

# ────────────────────── Save ──────────────────────

prs.save(str(OUTPUT_FILE))
print(f"\nSaved: {OUTPUT_FILE}")
print("Done!")
